import shutil
import subprocess
import zipfile
import time
from pathlib import Path
from tempfile import TemporaryDirectory

from PIL import Image, ImageDraw, ImageFont, ImageEnhance
from PyPDF2 import PdfReader, PdfWriter
from utils.config import config

try:
    import cv2
    VIDEO_TOOLS_AVAILABLE = True
except ImportError:
    cv2 = None
    VIDEO_TOOLS_AVAILABLE = False

try:
    from pdf2docx import Converter
except ImportError:
    Converter = None

try:
    import fitz
    PDF_RENDERING_AVAILABLE = True
except ImportError:
    fitz = None
    PDF_RENDERING_AVAILABLE = False

try:
    from docx import Document as WordDocument
except ImportError:
    WordDocument = None

try:
    from openpyxl import Workbook
except ImportError:
    Workbook = None

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    Presentation = None
    Inches = None

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    from PIL import ImageFile
    ImageFile.LOAD_TRUNCATED_IMAGES = True
except:
    pass


def _resolve_tesseract_executable() -> str | None:
    candidates = [
        Path(r"C:\Program Files\Tesseract-OCR\tesseract.exe"),
        Path(r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"),
    ]

    resolved = shutil.which("tesseract")
    if resolved:
        return resolved

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    return None


def _resolve_ffmpeg_executable() -> str | None:
    candidates = [
        Path(r"C:\ffmpeg\bin\ffmpeg.exe"),
        Path(r"C:\Program Files\ffmpeg\bin\ffmpeg.exe"),
        Path(r"C:\Program Files (x86)\ffmpeg\bin\ffmpeg.exe"),
    ]

    resolved = shutil.which("ffmpeg")
    if resolved:
        return resolved

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    return None


def _resolve_libreoffice_executable() -> str | None:
    candidates = [
        Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
        Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    ]

    resolved = shutil.which("soffice") or shutil.which("libreoffice")
    if resolved:
        return resolved

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    return None


def _resolve_ghostscript_executable() -> str | None:
    candidates = [
        Path(r"C:\Program Files\gs\gs10.04.0\bin\gswin64c.exe"),
        Path(r"C:\Program Files\gs\gs10.03.1\bin\gswin64c.exe"),
        Path(r"C:\Program Files\gs\gs10.02.1\bin\gswin64c.exe"),
    ]

    resolved = shutil.which("gswin64c") or shutil.which("gs")
    if resolved:
        return resolved

    for candidate in candidates:
        if candidate.exists():
            return str(candidate)

    return None


def is_libreoffice_available() -> bool:
    return _resolve_libreoffice_executable() is not None


def is_ghostscript_available() -> bool:
    return _resolve_ghostscript_executable() is not None


def is_conversion_available(conversion_target: str) -> bool:
    if conversion_target in {
        "word_to_pdf",
        "powerpoint_to_pdf",
        "excel_to_pdf",
        "html_to_pdf",
    }:
        return is_libreoffice_available()

    if conversion_target == "pdf_to_pdfa":
        return is_ghostscript_available()

    return True


def _convert_with_libreoffice(source_path: Path, target_extension: str, output_dir: Path) -> Path:
    libreoffice_executable = _resolve_libreoffice_executable()
    if not libreoffice_executable:
        raise ValueError(
            "This conversion is unavailable on this deployment because LibreOffice is not installed."
        )

    command = [
        libreoffice_executable,
        "--headless",
        "--convert-to",
        target_extension,
        "--outdir",
        str(output_dir),
        str(source_path),
    ]
    subprocess.run(command, check=True, capture_output=True, timeout=300)
    output_path = output_dir / f"{source_path.stem}.{target_extension}"
    if not output_path.exists():
        raise ValueError("LibreOffice could not create the converted file.")
    return output_path


def extract_zip_archive(zip_path: Path, output_dir: Path) -> list[Path]:
    output_dir.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(zip_path, "r") as zip_ref:
        file_members = [member for member in zip_ref.infolist() if not member.is_dir()]
        if len(file_members) > config.ZIP_MAX_EXTRACTED_FILES:
            raise ValueError(
                f"This ZIP contains too many files. The limit is {config.ZIP_MAX_EXTRACTED_FILES} files."
            )

        total_uncompressed_size = sum(member.file_size for member in file_members)
        if total_uncompressed_size > config.ZIP_MAX_EXTRACTED_SIZE:
            limit_mb = round(config.ZIP_MAX_EXTRACTED_SIZE / (1024 * 1024))
            raise ValueError(
                f"This ZIP expands beyond the allowed limit of {limit_mb} MB after extraction."
            )

        for member in zip_ref.infolist():
            member_path = output_dir / member.filename
            resolved_output_dir = output_dir.resolve()
            resolved_member_path = member_path.resolve()
            if not str(resolved_member_path).startswith(str(resolved_output_dir)):
                raise ValueError("ZIP contains an unsafe file path")

        zip_ref.extractall(output_dir)

    extracted_files = [path for path in output_dir.rglob("*") if path.is_file()]
    return extracted_files


def compress_image_file(source_path: Path, quality: int = 65) -> Path:
    output_path = source_path.with_name(f"{source_path.stem}_compressed.jpg")

    with Image.open(source_path) as image:
        rgb_image = image.convert("RGB")
        rgb_image.save(output_path, format="JPEG", optimize=True, quality=quality)

    return output_path


def convert_image_file(source_path: Path, conversion_target: str) -> Path:
    conversion_target = conversion_target.lower()

    if conversion_target == "jpg_to_pdf":
        conversion_target = "image_to_pdf"

    if conversion_target == "word_to_pdf":
        if source_path.suffix.lower() not in {".doc", ".docx"}:
            raise ValueError("Please send a DOC or DOCX file for Word to PDF conversion.")
        return _convert_with_libreoffice(source_path, "pdf", source_path.parent)

    if conversion_target == "powerpoint_to_pdf":
        if source_path.suffix.lower() not in {".ppt", ".pptx"}:
            raise ValueError("Please send a PPT or PPTX file for PowerPoint to PDF conversion.")
        return _convert_with_libreoffice(source_path, "pdf", source_path.parent)

    if conversion_target == "excel_to_pdf":
        if source_path.suffix.lower() not in {".xls", ".xlsx"}:
            raise ValueError("Please send an XLS or XLSX file for Excel to PDF conversion.")
        return _convert_with_libreoffice(source_path, "pdf", source_path.parent)

    if conversion_target == "html_to_pdf":
        if source_path.suffix.lower() not in {".html", ".htm"}:
            raise ValueError("Please send an HTML or HTM file for HTML to PDF conversion.")
        return _convert_with_libreoffice(source_path, "pdf", source_path.parent)

    if conversion_target == "pdf_to_jpg":
        return convert_pdf_to_jpg(source_path)

    if conversion_target == "pdf_to_word":
        if source_path.suffix.lower() != ".pdf":
            raise ValueError("Please send a PDF file for PDF to Word conversion.")
        if Converter is None:
            raise ValueError(
                "PDF to Word is unavailable because the required pdf2docx package is not installed."
            )

        output_path = source_path.with_suffix(".docx")
        converter = Converter(str(source_path))
        try:
            converter.convert(str(output_path), start=0, end=None)
        finally:
            converter.close()
        return output_path

    if conversion_target == "pdf_to_powerpoint":
        return convert_pdf_to_powerpoint(source_path)

    if conversion_target == "pdf_to_excel":
        return convert_pdf_to_excel(source_path)

    if conversion_target == "pdf_to_pdfa":
        return convert_pdf_to_pdfa(source_path)

    with Image.open(source_path) as image:
        if conversion_target == "jpg_to_png":
            output_path = source_path.with_suffix(".png")
            image.convert("RGBA").save(output_path, format="PNG")
            return output_path

        if conversion_target == "png_to_jpg":
            output_path = source_path.with_suffix(".jpg")
            image.convert("RGB").save(output_path, format="JPEG", quality=90)
            return output_path

        if conversion_target == "image_to_pdf":
            output_path = source_path.with_suffix(".pdf")
            image.convert("RGB").save(output_path, format="PDF")
            reader = PdfReader(str(output_path))
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with output_path.open("wb") as pdf_handle:
                writer.write(pdf_handle)
            return output_path

    raise ValueError("Unsupported conversion selected.")


def convert_pdf_to_jpg(source_path: Path) -> Path:
    if source_path.suffix.lower() != ".pdf":
        raise ValueError("Please send a PDF file for PDF to JPG conversion.")
    if not PDF_RENDERING_AVAILABLE:
        raise ValueError("PDF to JPG is unavailable because PyMuPDF is not installed.")

    output_dir = source_path.parent / f"{source_path.stem}_jpg_pages"
    output_dir.mkdir(parents=True, exist_ok=True)

    document = fitz.open(source_path)
    image_paths: list[Path] = []
    try:
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = output_dir / f"{source_path.stem}_page_{index}.jpg"
            pixmap.save(str(image_path))
            image_paths.append(image_path)
    finally:
        document.close()

    if not image_paths:
        raise ValueError("The PDF does not contain any renderable pages.")

    if len(image_paths) == 1:
        return image_paths[0]

    return create_zip_from_files([str(path) for path in image_paths])


def convert_pdf_to_powerpoint(source_path: Path) -> Path:
    if source_path.suffix.lower() != ".pdf":
        raise ValueError("Please send a PDF file for PDF to PowerPoint conversion.")
    if not PDF_RENDERING_AVAILABLE:
        raise ValueError("PDF to PowerPoint is unavailable because PyMuPDF is not installed.")
    if Presentation is None or Inches is None:
        raise ValueError("PDF to PowerPoint is unavailable because python-pptx is not installed.")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    with TemporaryDirectory(dir=source_path.parent) as temp_dir_name:
        temp_dir = Path(temp_dir_name)
        document = fitz.open(source_path)
        try:
            for index, page in enumerate(document, start=1):
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                image_path = temp_dir / f"page_{index}.jpg"
                pixmap.save(str(image_path))
                slide = presentation.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(image_path),
                    0,
                    0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )
        finally:
            document.close()

    output_path = source_path.with_suffix(".pptx")
    presentation.save(str(output_path))
    return output_path


def convert_pdf_to_excel(source_path: Path) -> Path:
    if source_path.suffix.lower() != ".pdf":
        raise ValueError("Please send a PDF file for PDF to Excel conversion.")
    if not PDF_RENDERING_AVAILABLE:
        raise ValueError("PDF to Excel is unavailable because PyMuPDF is not installed.")
    if Workbook is None:
        raise ValueError("PDF to Excel is unavailable because openpyxl is not installed.")

    workbook = Workbook()
    default_sheet = workbook.active
    workbook.remove(default_sheet)

    document = fitz.open(source_path)
    try:
        for index, page in enumerate(document, start=1):
            sheet = workbook.create_sheet(title=f"Page_{index}")
            text = page.get_text("text").splitlines()
            for row_index, line in enumerate(text, start=1):
                sheet.cell(row=row_index, column=1, value=line)
            if not text:
                sheet.cell(row=1, column=1, value="No extractable text found on this page.")
    finally:
        document.close()

    output_path = source_path.with_suffix(".xlsx")
    workbook.save(str(output_path))
    return output_path


def convert_pdf_to_pdfa(source_path: Path) -> Path:
    if source_path.suffix.lower() != ".pdf":
        raise ValueError("Please send a PDF file for PDF to PDF/A conversion.")

    ghostscript_executable = _resolve_ghostscript_executable()
    if not ghostscript_executable:
        raise ValueError("PDF to PDF/A is unavailable on this deployment because Ghostscript is not installed.")

    output_path = source_path.with_name(f"{source_path.stem}_pdfa.pdf")
    command = [
        ghostscript_executable,
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dNOOUTERSAVE",
        "-sProcessColorModel=DeviceRGB",
        "-sDEVICE=pdfwrite",
        "-dPDFACompatibilityPolicy=1",
        f"-sOutputFile={output_path}",
        str(source_path),
    ]
    subprocess.run(command, check=True, capture_output=True, timeout=300)
    if not output_path.exists():
        raise ValueError("Ghostscript could not create the PDF/A file.")
    return output_path


def rename_file_copy(source_path: Path, new_name: str) -> Path:
    output_path = source_path.with_name(new_name)
    shutil.copy2(source_path, output_path)
    return output_path


def merge_pdf_files(file_paths: list[str]) -> Path:
    if len(file_paths) < 2:
        raise ValueError("Please upload at least two PDF files to merge")

    source_paths = [Path(file_path) for file_path in file_paths]
    output_path = source_paths[0].parent / "merged_document.pdf"
    writer = PdfWriter()

    for pdf_path in source_paths:
        if pdf_path.suffix.lower() != ".pdf":
            raise ValueError("Only PDF files can be merged")
        reader = PdfReader(str(pdf_path))
        for page in reader.pages:
            writer.add_page(page)

    with output_path.open("wb") as file_handle:
        writer.write(file_handle)

    return output_path


def create_zip_from_files(file_paths: list[str]) -> Path:
    source_paths = [Path(file_path) for file_path in file_paths]
    if not source_paths:
        raise ValueError("Please upload at least one file to create a ZIP archive")

    output_path = source_paths[0].parent / "file_flex_bundle.zip"
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as zip_handle:
        seen_names: dict[str, int] = {}
        for file_path in source_paths:
            if not file_path.is_file():
                continue
            archive_name = file_path.name
            if archive_name in seen_names:
                seen_names[archive_name] += 1
                stem = file_path.stem
                suffix = file_path.suffix
                archive_name = f"{stem}_{seen_names[file_path.name]}{suffix}"
            else:
                seen_names[archive_name] = 1
            zip_handle.write(file_path, arcname=archive_name)

    return output_path


def enhance_image(source_path: Path, enhancement_type: str = "brightness") -> Path:
    with Image.open(source_path) as image:
        enhancer_map = {
            "brightness": ImageEnhance.Brightness,
            "contrast": ImageEnhance.Contrast,
            "sharpness": ImageEnhance.Sharpness,
            "color": ImageEnhance.Color,
        }
        
        enhancer_class = enhancer_map.get(enhancement_type, ImageEnhance.Brightness)
        enhancer = enhancer_class(image)
        enhanced = enhancer.enhance(1.3)
        
        output_path = source_path.with_name(f"{source_path.stem}_enhanced.jpg")
        enhanced.convert("RGB").save(output_path, format="JPEG", quality=90)
        return output_path


def resize_image(source_path: Path, width: int = 800, height: int = 600) -> Path:
    with Image.open(source_path) as image:
        image.thumbnail((width, height), Image.Resampling.LANCZOS)
        output_path = source_path.with_name(f"{source_path.stem}_resized.jpg")
        image.convert("RGB").save(output_path, format="JPEG", quality=90)
        return output_path


def extract_text_from_image(source_path: Path) -> str:
    if not TESSERACT_AVAILABLE:
        raise ValueError("OCR feature is not available. Install pytesseract and tesseract-ocr.")

    tesseract_executable = _resolve_tesseract_executable()
    if not tesseract_executable:
        raise ValueError(
            "OCR needs the Tesseract Windows app installed. Install Tesseract OCR and try again"
        )

    pytesseract.pytesseract.tesseract_cmd = tesseract_executable

    try:
        image = Image.open(source_path)
        text = pytesseract.image_to_string(image, lang="eng+amh")
        return text
    except Exception as e:
        raise ValueError(f"Failed to extract text from image: {str(e)}")


def add_watermark(source_path: Path, watermark_text: str = "File Flex Bot") -> Path:
    with Image.open(source_path) as image:
        watermarked = image.convert("RGB")
        draw = ImageDraw.Draw(watermarked, "RGBA")
        
        width, height = watermarked.size
        font_size = max(20, width // 15)
        
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            font = ImageFont.load_default()
        
        bbox = draw.textbbox((0, 0), watermark_text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        position = (width - text_width - 20, height - text_height - 20)
        draw.text(position, watermark_text, font=font, fill=(255, 255, 255, 128))
        
        output_path = source_path.with_name(f"{source_path.stem}_watermarked.jpg")
        watermarked.save(output_path, format="JPEG", quality=90)
        return output_path


def video_to_gif(source_path: Path, fps: int = 10, max_frames: int = 30) -> Path:
    if not VIDEO_TOOLS_AVAILABLE:
        raise ValueError("Video to GIF requires opencv-python to be installed")

    try:
        cap = cv2.VideoCapture(str(source_path))
        
        if not cap.isOpened():
            raise ValueError("Failed to open video file")
        
        frames = []
        frame_count = 0
        frame_interval = max(1, int(cap.get(cv2.CAP_PROP_FPS) / fps))
        
        while len(frames) < max_frames:
            ret, frame = cap.read()
            if not ret:
                break
            
            if frame_count % frame_interval == 0:
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                frame_pil = Image.fromarray(frame_rgb)
                frame_pil.thumbnail((400, 400), Image.Resampling.LANCZOS)
                frames.append(frame_pil)
            
            frame_count += 1
        
        cap.release()
        
        if not frames:
            raise ValueError("No frames extracted from video")
        
        output_path = source_path.with_suffix(".gif")
        frames[0].save(
            output_path,
            save_all=True,
            append_images=frames[1:],
            duration=100,
            loop=0,
            optimize=True
        )
        
        return output_path
    except Exception as e:
        raise ValueError(f"Failed to convert video to GIF: {str(e)}")


def compress_video(source_path: Path, quality: str = "medium") -> Path:
    quality_settings = {
        "low": (18, 800),
        "medium": (20, 1280),
        "high": (23, 1920),
    }
    
    crf, width = quality_settings.get(quality, quality_settings["medium"])
    
    output_path = source_path.with_name(f"{source_path.stem}_compressed.mp4")
    ffmpeg_executable = _resolve_ffmpeg_executable()
    if not ffmpeg_executable:
        raise ValueError("Video compression failed because ffmpeg is not installed")
    
    try:
        import subprocess
        command = [
            ffmpeg_executable, "-i", str(source_path),
            "-c:v", "libx264", "-crf", str(crf),
            "-vf", f"scale=min({width}\\,iw):-1",
            "-c:a", "aac", "-b:a", "128k",
            "-y", str(output_path)
        ]
        subprocess.run(command, check=True, capture_output=True, timeout=300)
        return output_path
    except Exception as e:
        raise ValueError(f"Video compression failed. Ensure ffmpeg is installed: {str(e)}")


def split_pdf(source_path: Path, start_page: int = 1, end_page: int = -1) -> Path:
    if source_path.suffix.lower() != ".pdf":
        raise ValueError("Please provide a PDF file")
    
    reader = PdfReader(str(source_path))
    total_pages = len(reader.pages)
    
    if start_page < 1 or (end_page != -1 and end_page > total_pages):
        raise ValueError(f"Invalid page range. PDF has {total_pages} pages.")
    
    writer = PdfWriter()
    for page_num in range(start_page - 1, min(end_page, total_pages)):
        writer.add_page(reader.pages[page_num])
    
    output_path = source_path.with_name(f"{source_path.stem}_pages_{start_page}-{end_page if end_page != -1 else total_pages}.pdf")
    with output_path.open("wb") as f:
        writer.write(f)
    
    return output_path


def get_file_info(file_path: Path) -> dict:
    stat = file_path.stat()
    return {
        "size": stat.st_size,
        "size_mb": round(stat.st_size / (1024 * 1024), 2),
        "created": stat.st_ctime,
        "modified": stat.st_mtime,
    }
